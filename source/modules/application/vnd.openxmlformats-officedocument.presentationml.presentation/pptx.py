# Copyright (C) 2013 Hogeschool van Amsterdam

# This program is free software; you can redistribute it and/or
# modify it under the terms of the GNU General Public License
# as published by the Free Software Foundation; either version 2
# of the License, or (at your option) any later version.

# This program is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
# GNU General Public License for more details.

#!/usr/bin/env python

# TABLE: title:LONGTEXT, created_by:LONGTEXT, modified_by:LONGTEXT, revision:INT, creation_date:LONGTEXT, modified_date:LONGTEXT, total_length: FLOAT, total_words:INT, application:LONGTEXT, pp_format:LONGTEXT, paragraphs:INT, slides:INT, notes:INT, hidden_slides:INT, videos:INT, company:LONGTEXT, shared:LONGTEXT, version:FLOAT, content:LONGTEXT

import xml.etree.ElementTree as ET
import re
import zipfile
import sys
from pptx import Presentation


def process(fullpath, config, rcontext, columns=None):
    try:
        document = zipfile.ZipFile(fullpath)
    except:
        exit()

    # document core, ie. keywords/title/subject and mod+creation dates
    xmlprop = document.read("docProps/core.xml")

    # data regarding ammount of pages/words. Also contains app version and OS.
    xmlapp = document.read("docProps/app.xml")

    # Minidom alternative
    tree = ET.fromstring(xmlprop)
    tree_app = ET.fromstring(xmlapp)

    #### Just some data, might be usefull lat0r ####
    data_tree = []
    data_app = []

# 	for basic in range(6):
# 		data_tree.append(tree[basic].text)

# 	for app in range(17):
# 		data_app.append(tree_app[app].text)


    data_tree.append(tree[0].text)
    data_tree.append(tree[1].text)
    data_tree.append(tree[2].text)
    data_tree.append(tree[3].text)
    data_tree.append(tree[4].text)
    data_tree.append(tree[5].text)

    data_app.append(tree_app[0].text)
    data_app.append(tree_app[1].text)
    data_app.append(tree_app[2].text)
    data_app.append(tree_app[3].text)
    data_app.append(tree_app[4].text)
    data_app.append(tree_app[5].text)
    data_app.append(tree_app[6].text)
    data_app.append(tree_app[7].text)
    data_app.append(tree_app[8].text)
    data_app.append(tree_app[12].text)
    data_app.append(tree_app[14].text)
    data_app.append(tree_app[15].text)


    textlist = []
    powerpoints = Presentation(fullpath)
    slidenum = 1
    for slide in powerpoints.slides:
        for shape in slide.shapes:
            if not shape.has_textframe:
                continue
            for paragraph in shape.textframe.paragraphs:
                for run in paragraph.runs:
                    if not any(d.get('Slide', None) == slidenum for d in textlist):
                        textlist.append({"Slide": slidenum, "Content": [run.text]})
                    else:
                        for val in textlist:
                            if val["Slide"] == slidenum:
                                val["Content"].append(run.text)
        slidenum += 1

    merged = data_tree + data_app + textlist
    return merged


# 	title = tree[0].text
# 	created = tree[1].text
# 	lastmod = tree[2].text
# 	revision = tree[3].text
# 	madeon = tree[4].text
# 	changedon = tree[5].text

# 	time = tree_app[0].text
# 	words = tree_app[1].text
# 	application = tree_app[2].text
# 	ppFormat = tree_app[3].text
# 	paragraphs = tree_app[4].text
# 	slides = tree_app[5].text
# 	notes = tree_app[6].text
# 	hiddenslides = tree_app[7].text
# 	multimediaclips = tree_app[8].text
# 	company = tree_app[12].text
# 	shared = tree_app[14].text
# 	version = tree_app[16].text
